"""[routes/reindex] 사내 문서 전체 재인덱싱 도메인 — 트리거·상태·실행·문서추출. APIRouter.
config·core·rag.manager·state(chunk_config)·handlers.files·deps 의존. 파서 라이브러리 지연import.
※ 리팩토링 이동 — 본문 byte-동일(@app→@router). run_reindex/_extract_file_for_reindex는 헬퍼."""
import os
import pickle
import asyncio
import traceback
import logging

from fastapi import APIRouter, Depends, HTTPException, Request, BackgroundTasks

from config import DATA_ROOT, DOCS_PATH
from core.retriever import HybridRetriever, KiwiBM25Retriever
from core.vector_store import ChromaVectorStore
from rag.manager import RAGManager
from state import chunk_config
from handlers.files import _extract_html_text
from deps import limiter, _require_admin

logger = logging.getLogger("main")

router = APIRouter()

# 재인덱싱 진행 상태(단일 작업). run_reindex가 in-place 갱신.
reindex_status: dict = {"status": "idle", "progress": "", "error": ""}


def _extract_file_for_reindex(fpath: str, fname: str) -> list:
    """확장자별 텍스트 추출 (동기, 스레드 풀용)"""
    from langchain_core.documents import Document as LCDoc
    ext = os.path.splitext(fname)[1].lower()
    try:
        if ext == ".pdf":
            import fitz
            docs, doc = [], fitz.open(fpath)
            for i, page in enumerate(doc):
                text = page.get_text().strip()
                if text:
                    docs.append(LCDoc(page_content=text, metadata={"source_file": fname, "page_no": i + 1}))
            doc.close()
            return docs
        elif ext in (".txt", ".md"):
            with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read().strip()
            return [LCDoc(page_content=text, metadata={"source_file": fname})] if text else []
        elif ext in (".html", ".htm"):
            with open(fpath, "rb") as f:
                text = _extract_html_text(f.read())
            return [LCDoc(page_content=text, metadata={"source_file": fname})] if text else []
        elif ext == ".docx":
            from docx import Document as DocxDoc
            from docx.table import Table
            from docx.text.paragraph import Paragraph
            doc = DocxDoc(fpath)
            # 문단과 표를 문서 순서대로 읽음(기존엔 paragraphs만 읽어 표 내용 누락).
            parts = []
            for child in doc.element.body.iterchildren():
                tag = child.tag
                if tag.endswith("}p"):
                    p = Paragraph(child, doc)
                    if p.text.strip():
                        parts.append(p.text.strip())
                elif tag.endswith("}tbl"):
                    for row in Table(child, doc).rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
            text = "\n".join(parts)
            return [LCDoc(page_content=text, metadata={"source_file": fname})] if text else []
        elif ext in (".xlsx", ".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(fpath, read_only=True, data_only=True)
            parts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                rows = [" | ".join(str(c) for c in row if c is not None) for row in ws.iter_rows(values_only=True)]
                rows = [r for r in rows if r.strip()]
                if rows:
                    parts.append(f"[시트: {sheet}]\n" + "\n".join(rows))
            wb.close()
            text = "\n\n".join(parts)
            return [LCDoc(page_content=text, metadata={"source_file": fname})] if text else []
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(fpath)
            parts = []
            for i, slide in enumerate(prs.slides):
                lines = []
                for s in slide.shapes:
                    # 표 도형은 s.text로 안 잡힘 → has_table로 별도 처리(기존 누락 보강)
                    if getattr(s, "has_table", False):
                        for row in s.table.rows:
                            cells = [c.text.strip() for c in row.cells if c.text.strip()]
                            if cells:
                                lines.append(" | ".join(cells))
                    elif hasattr(s, "text") and s.text.strip():
                        lines.append(s.text.strip())
                if lines:
                    parts.append(f"[슬라이드 {i + 1}]\n" + "\n".join(lines))
            text = "\n\n".join(parts)
            return [LCDoc(page_content=text, metadata={"source_file": fname})] if text else []
    except Exception as e:
        logger.warning(f"⚠️ [재인덱싱] '{fname}' 파싱 실패: {e}")
    return []

async def run_reindex():
    """사내 문서 DB 전체 재구축 백그라운드 태스크"""
    global reindex_status
    try:
        from core.text_splitter import ParentChildSplitter

        allowed_ext = {".pdf", ".txt", ".md", ".html", ".htm", ".docx", ".xlsx", ".xls", ".pptx"}
        files = [f for f in os.listdir(DOCS_PATH)
                 if os.path.isfile(os.path.join(DOCS_PATH, f))
                 and os.path.splitext(f)[1].lower() in allowed_ext]

        if not files:
            reindex_status.update({"status": "done", "progress": "인덱싱할 문서가 없습니다.", "error": ""})
            return

        # 1. 텍스트 추출 (LlamaParse 캐시 우선, 없으면 Cloud 파싱, 최종 fallback PyMuPDF)
        from core.document_loader import DocumentLoader
        _loader = DocumentLoader(
            document_dir=DOCS_PATH,
            use_marker=False,
            cache_dir=os.path.join(DATA_ROOT, "data/parsed_cache"),
        )
        all_docs = []
        for i, fname in enumerate(files):
            reindex_status["progress"] = f"파싱 중 ({i + 1}/{len(files)}): {fname}"
            fpath = os.path.join(DOCS_PATH, fname)
            if fname.lower().endswith(".pdf"):
                docs = await asyncio.to_thread(_loader.load_pdf, fpath)
            else:
                docs = await asyncio.to_thread(_extract_file_for_reindex, fpath, fname)
            all_docs.extend(docs)
            logger.info(f"  ↳ 📄 [재인덱싱] {fname}: {len(docs)}페이지")

        if not all_docs:
            reindex_status.update({"status": "error", "progress": "오류", "error": "모든 파일에서 텍스트 추출 실패"})
            return

        # 2. Parent-Child 분할 (chunk_config 적용 — 벤치마크 시 크기 조정 가능)
        reindex_status["progress"] = f"청킹 중... ({len(all_docs)}페이지)"
        splitter = ParentChildSplitter(
            parent_chunk_size=chunk_config["parent_size"],
            child_chunk_size=chunk_config["child_size"],
            child_chunk_overlap=chunk_config["child_overlap"],
            parent_chunk_overlap=chunk_config["parent_overlap"],
        )
        parent_map, child_docs = await asyncio.to_thread(splitter.split, all_docs)

        # 3. ChromaDB 전체 재구성
        reindex_status["progress"] = f"벡터 임베딩 중... ({len(child_docs)}개 청크)"
        vector_store = ChromaVectorStore(
            persist_directory=RAGManager.config["persist_dir"],
            collection_name="rag_documents",
            embeddings=RAGManager.embedder.get_embeddings()
        )
        await asyncio.to_thread(vector_store.create_from_documents, child_docs, True)

        # 4. BM25 재구성 + 디스크 저장
        reindex_status["progress"] = "BM25 인덱싱 중..."
        db_path = os.path.join(DATA_ROOT, "databases")
        parent_docs = list(parent_map.values())
        bm25_instance = await asyncio.to_thread(KiwiBM25Retriever, parent_docs)
        await asyncio.to_thread(bm25_instance.save_to_disk, os.path.join(db_path, "bm25_retriever.pkl"))

        with open(os.path.join(db_path, "parent_map.pkl"), "wb") as f:
            pickle.dump(parent_map, f)

        # 5. RAGManager 리트리버 무중단 교체
        reindex_status["progress"] = "RAG 엔진 갱신 중..."
        RAGManager.chain.retriever = HybridRetriever(
            parent_map=parent_map,
            child_vector_store=vector_store,
            parent_docs=parent_docs,
            bm25_instance=bm25_instance,
            top_k=RAGManager.top_k,
        )

        msg = f"완료! 문서 {len(files)}개 → parent {len(parent_map)}개 / child {len(child_docs)}개"
        reindex_status.update({"status": "done", "progress": msg, "error": ""})
        logger.info(f"✅ [재인덱싱] {msg}")

    except Exception as e:
        logger.error(f"❌ [재인덱싱] 실패: {traceback.format_exc()}")
        reindex_status.update({"status": "error", "progress": "오류 발생", "error": str(e)})

@router.post("/api/admin/reindex")
@limiter.limit("2/hour")
async def trigger_reindex(request: Request, background_tasks: BackgroundTasks, user_info: dict = Depends(_require_admin)):
    """사내 문서 전체 재인덱싱 트리거"""
    if reindex_status["status"] == "running":
        raise HTTPException(status_code=409, detail="재인덱싱이 이미 진행 중입니다.")
    reindex_status.update({"status": "running", "progress": "시작 중...", "error": ""})
    background_tasks.add_task(run_reindex)
    return {"success": True, "message": "재인덱싱이 백그라운드에서 시작되었습니다."}

@router.get("/api/admin/reindex/status")
async def get_reindex_status(user_info: dict = Depends(_require_admin)):
    """재인덱싱 진행 상태 조회"""
    return reindex_status
