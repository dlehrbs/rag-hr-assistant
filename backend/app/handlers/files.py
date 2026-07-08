"""[handlers/files] 파일 파싱 토대 — PDF(fitz)/스캔OCR폴백(LlamaParse)/HTML/DOCX/PPTX/XLSX 추출 +
파싱결과 캐시(content 해시). 파서 라이브러리는 함수 내부 지연 import(선택적 의존).
config·stdlib에만 의존(순환0). main(temp-upload·reindex)·projects_logic 공유.
※ 리팩토링으로 main.py에서 이동 — 함수 본문 byte-동일."""
import os
import io
import json
import asyncio
import hashlib
import tempfile
import logging

from config import PROJECT_PARSE_CACHE, _SCAN_OCR_FALLBACK, _SCAN_CHAR_THRESHOLD

logger = logging.getLogger("main")


def _extract_pdf_sync(content: bytes, filename: str) -> list:
    """[C-04] 별도 스레드에서 실행되는 동기 PDF 파싱 함수.
    fitz.open()은 CPU 집중 작업이므로 이벤트 루프 차단 방지를 위해 분리."""
    import fitz
    from langchain_core.documents import Document
    documents = []
    doc = None
    try:
        doc = fitz.open(stream=content, filetype="pdf")
        # [E-03] 암호화 PDF 방어
        if doc.is_encrypted:
            raise ValueError("암호화된 PDF는 지원하지 않습니다. 잠금 해제 후 재업로드해주세요.")
        for i, page in enumerate(doc):
            try:
                text = page.get_text().strip()
                if text:
                    documents.append(Document(
                        page_content=text,
                        metadata={"page_no": i + 1, "source": filename}
                    ))
            except Exception:
                continue  # 특정 페이지 파싱 실패해도 나머지 계속 진행
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"PDF 파싱 실패: {e}")
    finally:
        if doc:
            doc.close() # [C-04] 파일 디스크립터 누수 방지
    return documents

def _extract_html_text(content: bytes) -> str:
    """HTML 바이트 → 본문 텍스트. <script>/<style>/<head> 제거, <table>은 행별 ' | '로
    구조를 보존(office 파서와 동일한 표 표기 규칙). lxml은 python-docx의 필수 의존성이라 항상 존재."""
    from lxml import html as lxml_html
    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError:
        text = content.decode("cp949", errors="ignore")
    if not text.strip():
        return ""
    try:
        tree = lxml_html.fromstring(text)
    except Exception:
        return ""
    # 노이즈 요소 제거(스크립트·스타일·헤드 등)
    for bad in tree.xpath("//script | //style | //noscript | //head"):
        parent = bad.getparent()
        if parent is not None:
            parent.remove(bad)
    parts = []
    # 표는 구조 보존을 위해 별도 추출 후 트리에서 제거(본문 중복 방지)
    for table in tree.xpath("//table"):
        for row in table.xpath(".//tr"):
            cells = [" ".join(c.text_content().split()) for c in row.xpath("./td | ./th")]
            cells = [c for c in cells if c]
            if cells:
                parts.append(" | ".join(cells))
        parent = table.getparent()
        if parent is not None:
            parent.remove(table)
    # 나머지 본문 텍스트(빈 줄 정리)
    body = "\n".join(ln.strip() for ln in tree.text_content().splitlines() if ln.strip())
    if parts:
        body = (body + "\n" + "\n".join(parts)).strip()
    return body.strip()

def _parse_cache_path(content: bytes) -> str:
    import hashlib
    return os.path.join(PROJECT_PARSE_CACHE, f"{hashlib.sha256(content).hexdigest()}.json")

def _load_parse_cache(content: bytes):
    """캐시 히트 시 [{page_no, text}] 반환. 없거나 손상·비어있으면 None(미스 취급)."""
    path = _parse_cache_path(content)
    if not os.path.exists(path):
        return None
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        if isinstance(data, list) and any((p.get("text") or "").strip() for p in data):
            return data
    except Exception as e:
        logger.warning(f"파싱 캐시 읽기 실패(미스 처리): {e}")
    return None

def _save_parse_cache(content: bytes, pages: list):
    """원자적 저장(임시파일→rename)으로 반쪽 파일 방지. 실패는 무시(캐시는 보너스)."""
    try:
        os.makedirs(PROJECT_PARSE_CACHE, exist_ok=True)
        path = _parse_cache_path(content)
        tmp = f"{path}.tmp{os.getpid()}"
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(pages, f, ensure_ascii=False)
        os.replace(tmp, path)
    except Exception as e:
        logger.warning(f"파싱 캐시 저장 실패(무시): {e}")

def _pdf_page_count(content: bytes) -> int:
    """PDF 총 페이지 수(스캔 판정용). 실패 시 0."""
    try:
        import fitz
        d = fitz.open(stream=content, filetype="pdf")
        n = d.page_count
        d.close()
        return n
    except Exception:
        return 0

async def _parse_pdf_llamaparse(content: bytes, filename: str):
    """LlamaParse 정밀(markdown) 파싱 — 캐시 우선. 성공 시 Document 리스트, 실패/빈결과/키없음 시 None.
    정밀 모드와 '스캔본 자동 폴백' 양쪽에서 재사용."""
    from langchain_core.documents import Document
    api_key = os.getenv("LLAMA_CLOUD_API_KEY")
    if not api_key:
        logger.warning(f"  ↳ ⚠️ [정밀] LLAMA_CLOUD_API_KEY 미설정: '{filename}'")
        return None
    # 0) 캐시 적중 확인 (같은 내용이면 LlamaParse 미호출 = 한도 0)
    cached = _load_parse_cache(content)
    if cached:
        logger.info(f"  ↳ ♻️ [정밀] 캐시 적중 → LlamaParse 미호출: '{filename}' ({len(cached)}p)")
        return [Document(page_content=p["text"], metadata={"page_no": p.get("page_no", i + 1), "source": filename})
                for i, p in enumerate(cached)]
    try:
        from llama_parse import LlamaParse
        parser = LlamaParse(result_type="markdown", language="ko", verbose=False)
        with tempfile.NamedTemporaryFile(suffix=f"_{filename}", delete=True) as tmp:
            tmp.write(content)
            tmp.flush()
            llama_docs = await parser.aload_data(tmp.name)
        pages = [{"page_no": i + 1, "text": d.text} for i, d in enumerate(llama_docs)]
        docs = [Document(page_content=p["text"], metadata={"page_no": p["page_no"], "source": filename})
                for p in pages]
        if docs and any(d.page_content.strip() for d in docs):
            _save_parse_cache(content, pages)   # 성공·비어있지 않을 때만 캐시 저장
            logger.info(f"  ↳ ✅ [정밀] LlamaParse 파싱 성공·캐시 저장: '{filename}' ({len(docs)}p)")
            return docs
        logger.warning(f"  ↳ ⚠️ [정밀] LlamaParse 결과 비어있음: '{filename}'")
    except Exception as e:
        logger.warning(f"  ↳ ⚠️ [정밀] LlamaParse 실패(한도/오류): '{filename}' ({e})")
    return None

async def _parse_upload_to_documents(content: bytes, filename: str, mode: str) -> list:
    """업로드 파일(PDF/txt/md/docx/xlsx/pptx) → langchain Document 리스트로 파싱.
    PDF는 기존 임시업로드 경로(PyMuPDF / LlamaParse)를 재사용, 나머지는 직접 추출(BytesIO)."""
    import io
    from langchain_core.documents import Document
    ext = os.path.splitext(filename)[1].lower()

    if ext in (".txt", ".md"):
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("cp949", errors="ignore")
        if not text.strip():
            return []
        return [Document(page_content=text, metadata={"page_no": 1, "source": filename})]

    if ext in (".html", ".htm"):
        text = await asyncio.to_thread(_extract_html_text, content)
        return [Document(page_content=text, metadata={"page_no": 1, "source": filename})] if text.strip() else []

    # ── Office 문서 (동기 추출이라 스레드 풀로 위임) ──
    if ext == ".docx":
        def _docx():
            from docx import Document as DocxDoc
            from docx.table import Table
            from docx.text.paragraph import Paragraph
            doc = DocxDoc(io.BytesIO(content))
            # 문단과 표를 '문서 순서대로' 읽어 표가 주변 문맥과 분리되지 않게 함.
            # (기존엔 doc.paragraphs만 읽어 표 내용 전체가 누락됐음 — 표 든 보고서에서 치명적)
            parts = []
            for child in doc.element.body.iterchildren():
                tag = child.tag
                if tag.endswith("}p"):
                    p = Paragraph(child, doc)
                    if p.text.strip():
                        parts.append(p.text.strip())
                elif tag.endswith("}tbl"):
                    for row in Table(child, doc).rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
            return "\n".join(parts)
        text = await asyncio.to_thread(_docx)
        return [Document(page_content=text, metadata={"page_no": 1, "source": filename})] if text.strip() else []

    if ext in (".xlsx", ".xls"):
        def _xlsx():
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            parts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                rows = [" | ".join(str(c) for c in row if c is not None) for row in ws.iter_rows(values_only=True)]
                rows = [r for r in rows if r.strip()]
                if rows:
                    parts.append(f"[시트: {sheet}]\n" + "\n".join(rows))
            wb.close()
            return "\n\n".join(parts)
        text = await asyncio.to_thread(_xlsx)
        return [Document(page_content=text, metadata={"page_no": 1, "source": filename})] if text.strip() else []

    if ext == ".pptx":
        def _pptx():
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            parts = []
            for i, slide in enumerate(prs.slides):
                lines = []
                for s in slide.shapes:
                    # 표 도형은 s.text로 안 잡힘 → has_table로 별도 처리(기존 누락 보강)
                    if getattr(s, "has_table", False):
                        for row in s.table.rows:
                            cells = [c.text.strip() for c in row.cells if c.text.strip()]
                            if cells:
                                lines.append(" | ".join(cells))
                    elif hasattr(s, "text") and s.text.strip():
                        lines.append(s.text.strip())
                if lines:
                    parts.append(f"[슬라이드 {i + 1}]\n" + "\n".join(lines))
            return "\n\n".join(parts)
        text = await asyncio.to_thread(_pptx)
        return [Document(page_content=text, metadata={"page_no": 1, "source": filename})] if text.strip() else []

    if ext == ".pdf":
        if mode == "quality":
            # 정밀(LlamaParse). 실패(키없음/한도/네트워크)면 아래 고속(PyMuPDF)으로 자동 폴백.
            docs = await _parse_pdf_llamaparse(content, filename)
            if docs:
                return docs
            logger.warning(f"  ↳ ⚠️ [정밀] → 일반(PyMuPDF) 폴백: '{filename}'")
        # fast (PyMuPDF) — 기본값이자 정밀 실패 시 폴백
        docs = await asyncio.to_thread(_extract_pdf_sync, content, filename)
        # ── 스캔본(이미지 PDF) 자동 감지 → LlamaParse 폴백 ──
        #   고속 추출 글자수가 페이지당 임계 미만이면 "글자가 픽셀로 갇힌" 스캔본으로 보고
        #   OCR 내장 LlamaParse로 자동 재시도. (정밀 모드는 이미 위에서 시도했으므로 제외)
        if _SCAN_OCR_FALLBACK and mode != "quality":
            n_pages = _pdf_page_count(content)
            total_chars = sum(len(d.page_content.strip()) for d in docs)
            if n_pages > 0 and total_chars < n_pages * _SCAN_CHAR_THRESHOLD:
                logger.info(f"  ↳ 🔍 [스캔감지] '{filename}' 추출 {total_chars}자/{n_pages}p "
                            f"(임계 {_SCAN_CHAR_THRESHOLD}/p 미만) → LlamaParse 자동 폴백")
                lp = await _parse_pdf_llamaparse(content, filename)
                if lp:
                    return lp
                logger.warning(f"  ↳ ⚠️ [스캔감지] LlamaParse 폴백 실패 → PyMuPDF 결과 반환: '{filename}'")
        return docs

    raise ValueError(f"지원하지 않는 파일 형식입니다: {ext} (PDF/txt/md/html/docx/xlsx/pptx 가능)")
